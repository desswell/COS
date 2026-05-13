"""Quick QA: extract text from the generated presentation."""
from pathlib import Path
from pptx import Presentation

p = Path(__file__).parent / "presentation.pptx"
prs = Presentation(p)
out = []
out.append(f"Total slides: {len(prs.slides)}")
out.append(f"Slide size: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} in")
for i, slide in enumerate(prs.slides, 1):
    out.append(f"\n=== Slide {i} ===")
    n_text = 0
    n_image = 0
    n_chart = 0
    n_shape = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs)
                if txt.strip():
                    out.append(f"  T  {txt}")
                    n_text += 1
        elif shape.shape_type == 13:  # picture
            out.append(f"  IMG  x={shape.left/914400:.2f} y={shape.top/914400:.2f} w={shape.width/914400:.2f} h={shape.height/914400:.2f}")
            n_image += 1
        elif hasattr(shape, "chart") and shape.has_chart:
            out.append(f"  CHART  x={shape.left/914400:.2f} y={shape.top/914400:.2f} w={shape.width/914400:.2f} h={shape.height/914400:.2f}")
            n_chart += 1
        else:
            n_shape += 1
    out.append(f"  ... [text={n_text} img={n_image} chart={n_chart} shape={n_shape}]")

Path(__file__).parent.joinpath("qa.txt").write_text("\n".join(out), encoding="utf-8")
print("Wrote qa.txt")
